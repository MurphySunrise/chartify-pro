"""
PPT Report Generator Module
Creates PowerPoint reports with charts organized by significance.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from io import BytesIO
import tempfile
import os

from matplotlib.figure import Figure

from ..stats.calculator import DataTypeStats


class PPTGenerator:
    """Generates PowerPoint reports from analysis results."""
    
    # Layout settings
    CHARTS_PER_PAGE = 6  # 2 rows x 3 columns
    CHARTS_PER_ROW = 3
    
    # Slide dimensions (standard 16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    
    # Chart dimensions and positions
    MARGIN_LEFT = Inches(0.3)
    MARGIN_TOP = Inches(1.0)  # Leave room for title
    MARGIN_RIGHT = Inches(0.3)
    MARGIN_BOTTOM = Inches(0.3)
    CHART_SPACING = Inches(0.2)
    
    def __init__(self):
        """Initialize PPT generator."""
        self.prs: Optional[Presentation] = None
    
    def _calculate_chart_size(self) -> Tuple[float, float]:
        """Calculate chart width and height based on layout."""
        available_width = self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT
        available_height = self.SLIDE_HEIGHT - self.MARGIN_TOP - self.MARGIN_BOTTOM
        
        # Account for spacing
        total_h_spacing = self.CHART_SPACING * (self.CHARTS_PER_ROW - 1)
        total_v_spacing = self.CHART_SPACING * 1  # 2 rows = 1 gap
        
        chart_width = (available_width - total_h_spacing) / self.CHARTS_PER_ROW
        chart_height = (available_height - total_v_spacing) / 2  # 2 rows
        
        return chart_width, chart_height
    
    def _get_chart_position(self, index: int) -> Tuple[float, float]:
        """
        Get position for chart at given index (0-5).
        
        Args:
            index: Chart index on page (0-5)
            
        Returns:
            Tuple of (left, top) position
        """
        chart_width, chart_height = self._calculate_chart_size()
        
        row = index // self.CHARTS_PER_ROW
        col = index % self.CHARTS_PER_ROW
        
        left = self.MARGIN_LEFT + col * (chart_width + self.CHART_SPACING)
        top = self.MARGIN_TOP + row * (chart_height + self.CHART_SPACING)
        
        return left, top
    
    def create_presentation(self, template_path: Optional[str] = None) -> None:
        """
        Create a new presentation or load template.
        
        Args:
            template_path: Optional path to template file
        """
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            # Set 16:9 aspect ratio
            self.prs.slide_width = self.SLIDE_WIDTH
            self.prs.slide_height = self.SLIDE_HEIGHT
    
    def add_slide_with_title(self, title: str) -> 'Slide':
        """
        Add a new slide with title.
        
        Args:
            title: Slide title text
            
        Returns:
            The new slide object
        """
        # Use blank layout
        blank_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Add title text box
        left = Inches(0.5)
        top = Inches(0.2)
        width = self.SLIDE_WIDTH - Inches(1)
        height = Inches(0.6)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_chart_to_slide(
        self,
        slide,
        figure: Figure,
        index: int
    ) -> None:
        """
        Add a chart figure to slide at specified position.
        
        Args:
            slide: PowerPoint slide object
            figure: Matplotlib figure
            index: Position index (0-5)
        """
        # Save figure to temporary file
        chart_width, chart_height = self._calculate_chart_size()
        left, top = self._get_chart_position(index)
        
        # Convert figure to image bytes
        buf = BytesIO()
        figure.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        buf.seek(0)
        
        # Add picture to slide
        slide.shapes.add_picture(buf, left, top, width=chart_width, height=chart_height)
    
    def generate_report(
        self,
        figures: Dict[str, Figure],
        stats: Dict[str, DataTypeStats],
        csv_path: str,
        template_path: Optional[str] = None,
        progress_callback: Optional[callable] = None
    ) -> str:
        """
        Generate complete PPT report.
        
        Args:
            figures: Dictionary mapping data_type -> Figure
            stats: Dictionary mapping data_type -> DataTypeStats
            csv_path: Path to source CSV file (for output naming)
            template_path: Optional template file path
            progress_callback: Optional progress callback
            
        Returns:
            Path to generated PPT file
        """
        self.create_presentation(template_path)
        
        # Separate mismatch (p<0.05, significant) and match (p>=0.05, non-significant)
        # p < 0.05 means mismatch (significant difference from control)
        # p >= 0.05 means match (no significant difference from control)
        mismatch = []  # significant results (p < 0.05)
        match = []     # non-significant results (p >= 0.05)
        
        for data_type, stat in stats.items():
            if data_type in figures:
                if stat.has_significant_results():
                    mismatch.append(data_type)  # p < 0.05 -> mismatch
                else:
                    match.append(data_type)     # p >= 0.05 -> match
        
        # Sort each list
        mismatch.sort()
        match.sort()
        
        total_charts = len(figures)
        charts_processed = 0
        
        # Helper to add charts in batches
        def add_charts_batch(data_types: List[str], title: str):
            nonlocal charts_processed
            
            for i in range(0, len(data_types), self.CHARTS_PER_PAGE):
                batch = data_types[i:i + self.CHARTS_PER_PAGE]
                slide = self.add_slide_with_title(title)
                
                for j, data_type in enumerate(batch):
                    fig = figures[data_type]
                    self.add_chart_to_slide(slide, fig, j)
                    charts_processed += 1
                    
                    if progress_callback:
                        progress_callback(charts_processed / total_charts * 100)
        
        # Add mismatch charts first (p < 0.05, significant difference)
        if mismatch:
            add_charts_batch(mismatch, "metrics mismatch")
        
        # Then match charts (p >= 0.05, no significant difference)
        if match:
            add_charts_batch(match, "metrics match")
        
        # Generate output path
        csv_file = Path(csv_path)
        output_path = csv_file.parent / f"{csv_file.stem}_statistic_report.pptx"
        
        # Save presentation
        self.prs.save(str(output_path))
        
        return str(output_path)
    
    def close_figures(self, figures: Dict[str, Figure]) -> None:
        """
        Close all matplotlib figures to free memory.
        
        Args:
            figures: Dictionary of figures to close
        """
        import matplotlib.pyplot as plt
        for fig in figures.values():
            plt.close(fig)
