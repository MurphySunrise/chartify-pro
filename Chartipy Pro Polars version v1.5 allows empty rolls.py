import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
import polars as pl
from scipy.stats import ttest_ind
import scipy.stats as stats
from tkinter import Tk, filedialog, StringVar, DoubleVar, Button, Label, Frame, Radiobutton, Listbox, EXTENDED, messagebox, Entry, ttk
from tkinter.ttk import Combobox
import time
from io import BytesIO
import threading
from pptx import Presentation
from pptx.util import Inches
import webbrowser
import configparser
from concurrent.futures import ProcessPoolExecutor, as_completed

# 配置管理类
class ConfigManager:
    def __init__(self):
        self.config_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'chartify_config.ini')
        self.config = configparser.ConfigParser()
        
    def load_config(self):
        # 加载配置文件，返回PPT模板路径
        if os.path.exists(self.config_file):
            self.config.read(self.config_file)
            if 'Settings' in self.config and 'ppt_template' in self.config['Settings']:
                return self.config['Settings']['ppt_template']
        return None
        
    def save_config(self, ppt_template_path):
        # 保存PPT模板路径到配置文件
        if not self.config.has_section('Settings'):
            self.config.add_section('Settings')
        self.config['Settings']['ppt_template'] = ppt_template_path
        with open(self.config_file, 'w') as configfile:
            self.config.write(configfile)

def generate_item_sort(data, item_col, selected_data_cols, data_structure):
    if data_structure == "single":
        # 单列模式：从 item 列中提取唯一值并保持顺序
        items = data.select(item_col).to_numpy().flatten()
        unique_items, indices = np.unique(items, return_index=True)
        unique_items = unique_items[np.argsort(indices)]  # 根据首次出现的索引排序
    else:
        # 多列模式：使用选中的数据列名
        unique_items = selected_data_cols

    # 创建 item_sort 映射表
    item_sort_df = pl.DataFrame({
        "item": unique_items,
        "item_sort": range(1, len(unique_items) + 1)
    })

    print(item_sort_df)

    return item_sort_df

def process_data(data, group_col, item_col, selected_data_cols, data_structure, control_group, progress_var, status_label, root, csv_path):
    try:
        # 更新状态为开始计算
        status_label.config(text="Statistic Data Calculating...")
        root.update_idletasks()

        # 开始计时
        start_time = time.time()
        print("开始计算")

        if data_structure == "single":
            if len(selected_data_cols) != 1:
                raise ValueError("在单列结构下，必须且只能选择一个数据列。")
            long_data = data.select([group_col, item_col, selected_data_cols[0]]).rename({
                group_col: "group",
                item_col: "item",
                selected_data_cols[0]: "value"
            })
        else:
            long_data = data.unpivot(
                index=[group_col],
                on=selected_data_cols,
                variable_name="item",
                value_name="value"
            ).rename({
                group_col: "group"
            })
            long_data = long_data.filter(pl.col("value").is_not_null())

        # 计算统计结果
        
        # 补全所有 item 和 group 的组合
        if data_structure == "multiple":
            all_items = pl.DataFrame({"item": selected_data_cols})
        else:
            all_items = long_data.select("item").unique()

        all_groups = data.select(group_col).unique().rename({group_col: "group"})
        all_combinations = all_items.join(all_groups, how="cross")

        # 计算统计结果
        StaticResult = long_data.group_by(["item", "group"]).agg([
            pl.col("value").count().alias("count"),
            pl.col("value").mean().alias("mean"),
            pl.col("value").median().alias("median"),
            pl.col("value").quantile(0.05).alias("q5"),
            pl.col("value").quantile(0.95).alias("q95"),
            pl.col("value").std().alias("std"),
        ])

        # 合并并填补缺失值
        StaticResult = all_combinations.join(StaticResult, on=["item", "group"], how="left")
        StaticResult = StaticResult.with_columns([
            pl.col("count").fill_null(0),
        ])

        # 排序
        StaticResult = StaticResult.sort(["item", "group"])


        # 计算 Sigma Delta
        sigma_delta_StaticResults = []
        for row in StaticResult.iter_rows(named=True):
            if row["group"] == control_group:
                sigma_delta = float('nan')
            else:
                control_row = StaticResult.filter(
                    (pl.col("group") == control_group) & (pl.col("item") == row["item"])
                )
                if not control_row.is_empty():
                    control_mean = control_row["mean"][0]
                    control_std = control_row["std"][0]
                    sigma_delta = (row["mean"] - control_mean) / control_std if control_std != 0 else float('nan')
                else:
                    sigma_delta = float('nan')
            sigma_delta_StaticResults.append(sigma_delta)

        StaticResult = StaticResult.with_columns(pl.Series(sigma_delta_StaticResults).alias("sigma_delta"))

        # 转换 long_data 到 NumPy 格式
        item_data_dict = {}
        for item in long_data.select("item").unique().to_numpy().flatten():
            item_data = long_data.filter(pl.col("item") == item)
            groups = item_data.select("group").to_numpy().flatten()
            values = item_data.select("value").to_numpy().flatten()
            unique_groups = np.unique(groups)
            item_data_dict[item] = (groups, values, unique_groups)

        # 计算 p 值
        p_values = []
        total_items = len(StaticResult)
        for idx, row in enumerate(StaticResult.iter_rows(named=True)):
            item = row["item"]
            current_group = row["group"]
            if current_group == control_group:
                p_values.append(float('nan'))
                continue
            control_values = item_data_dict[item][1][item_data_dict[item][0] == control_group]
            current_values = item_data_dict[item][1][item_data_dict[item][0] == current_group]
            if len(control_values) < 2 or len(current_values) < 2:
                p_values.append(float('nan'))
                continue
            t_stat, p_val = ttest_ind(current_values, control_values, equal_var=False)
            p_values.append(p_val)

            # 更新进度条
            progress_var.set((idx + 1) / total_items * 50)  # 计算P值占总进度的50%
            root.update_idletasks()  # 强制刷新界面

        StaticResult = StaticResult.with_columns(pl.Series(p_values).alias("p_value"))

        # 结束计时并打印耗时
        end_time = time.time()
        print(f"Statistics Calculation Completed in {end_time - start_time:.2f} seconds")

        # 更新状态为计算完成
        status_label.config(text="Statistical Calculations Completed")
        root.update_idletasks()

        return StaticResult, item_data_dict
    except Exception as e:
        print(f"数据处理时出错: {e}")
        raise

def plot_item(item, StaticResult, control_group, item_data_dict, sample_size):
    try:
        # 生成单个项目的图表并返回图像数据
        groups, values, unique_groups = item_data_dict[item]
        colors = ['blue', 'red', 'green', 'purple', 'orange', 'cyan', 'magenta']

        # 如果某个item的行数超过指定的采样大小，进行采样
        sampled_groups = []
        sampled_values = []
        for g in unique_groups:
            group_mask = groups == g
            group_values = values[group_mask]
            group_indices = np.where(group_mask)[0]  # 获取当前组的索引
            if len(group_values) > sample_size:
                sample_indices = np.random.choice(len(group_values), sample_size, replace=False)
                sampled_groups.append(groups[group_indices][sample_indices])
                sampled_values.append(group_values[sample_indices])
            else:
                sampled_groups.append(groups[group_indices])
                sampled_values.append(group_values)

        # 合并采样后的数据
        groups = np.concatenate(sampled_groups)
        values = np.concatenate(sampled_values)

        # 使用 gridspec 来定义布局
        fig = plt.figure(figsize=(16, 10))
        gs = fig.add_gridspec(2, 3, height_ratios=[3, 1], width_ratios=[9, 9, 1])  # 修改为2x3布局

              
       # Box Chart
        ax_box = fig.add_subplot(gs[0, 0])
        box_data = [values[groups == g] for g in unique_groups]
        positions = np.arange(1, len(unique_groups) + 1)
        box = ax_box.boxplot(box_data, tick_labels=unique_groups, positions=positions, patch_artist=True)

        # 设置箱体颜色和透明度
        for patch, color in zip(box['boxes'], colors):
            patch.set_facecolor(color)
            patch.set_alpha(0.4)

        medians = [np.median(values[groups == g]) for g in unique_groups]

        # 在每个箱形图的中位数位置绘制一条横线
        for i, median in enumerate(medians):
            ax_box.plot([positions[i] - 0.2, positions[i] + 0.2], [median, median], color=colors[i % len(colors)], linewidth=4)

        ax_box.plot(positions, medians, marker='o', linestyle='-', color='blue')

        # 计算最大重复次数
        max_count = max(np.unique(values, return_counts=True)[1])

        # 添加动态递减抖动到散点图
        for i, g in enumerate(unique_groups):
            group_values = values[groups == g]
            unique_values, counts = np.unique(group_values, return_counts=True)
            for value, count in zip(unique_values, counts):
                # 根据最大重复次数动态调整抖动范围
                if max_count == 1:
                    jitter_range = 0
                else:
                    jitter_range = 0.5 * (count / max_count) # 动态调整抖动范围
                jitter = np.linspace(-jitter_range, jitter_range, count)  # 线性分布的抖动
                ax_box.scatter(np.full(count, positions[i]) + jitter, np.full(count, value), alpha=0.5, color=colors[i % len(colors)])

        ax_box.set_title(f'Box Chart for \n{item}', fontsize=30, fontweight='bold')
        ax_box.set_xlabel('\nGroup', fontsize=18, fontweight='bold')
        ax_box.set_ylabel('Value', fontsize=16, fontweight='bold')
        # 设置横坐标标签字体大小
        ax_box.tick_params(axis='x', labelsize=14)  # 调整这里的数值来改变横坐标标签大小
        for spine in ax_box.spines.values():
            spine.set_linewidth(2)

        # Normal Quantile Chart
        ax_quantile = fig.add_subplot(gs[0, 1])
        for i, g in enumerate(unique_groups):
            group_values = values[groups == g]
            (osm, osr), _ = stats.probplot(group_values, dist="norm")  # 将分位数转换为正态分位数
            ax_quantile.plot(osm, osr, marker='o', linestyle='-', label=g, color=colors[i % len(colors)])

        ax_quantile.set_title(f'Normal Quantile Chart for \n{item}', fontsize=30, fontweight='bold')
        ax_quantile.set_xlabel('Normal Quantile', fontsize=18, fontweight='bold')
        ax_quantile.tick_params(axis='x', labelsize=16)

        # 设置X轴刻度纵向显示
        quantiles = np.array([0.001, 0.005, 0.01, 0.02, 0.05, 0.1, 0.2, 0.3, 0.4, 0.5, 0.6, 0.7, 0.8, 0.9, 0.95, 0.99, 0.995, 0.999])
        # 使用 ppf 函数计算这些分位数在标准正态分布中的位置
        xticks = stats.norm.ppf(quantiles)
        ax_quantile.set_xticks(xticks)
        ax_quantile.set_xticklabels([str(q) for q in quantiles], rotation=90)
        ax_quantile.grid(axis='x', linestyle='--', color='gray', linewidth=0.5)
        for spine in ax_quantile.spines.values():
            spine.set_linewidth(2)

        # 空白图
        ax_blank = fig.add_subplot(gs[0, 2])  # 新增的空白图位置
        ax_blank.axis('off')  # 不显示坐标轴

        # 在空白图中添加图例
        handles, labels = ax_quantile.get_legend_handles_labels()  # 获取正态分位数图的图例项
        ax_blank.legend(handles, labels, loc='upper left', bbox_to_anchor=(-1.5, 1), fontsize=18)

        # Summary Table
        ax_table = fig.add_subplot(gs[1, :])
        item_StaticResult = StaticResult.filter(pl.col("item") == item)
        table_data = item_StaticResult.select(["group", "count", "mean", "median", "q5", "q95", "sigma_delta", "p_value"]).to_numpy()
        table_data[:, 1] = table_data[:, 1].astype(int)  # 确保count为整数
        numeric_indices = [2, 3, 4, 5, 6, 7]  # 调整索引，因为count现在是整数
        for i in numeric_indices:
            table_data[:, i] = np.round(table_data[:, i].astype(float), 4)

        table = ax_table.table(cellText=table_data, colLabels=["Group", "Count", "Mean", "Median", "Q5", "Q95", "Sigma_delta", "P_value"], loc='center', cellLoc='center')
        table.scale(1, 1.5)
        ax_table.axis('off')

        # 设置表格边框宽度和表头样式
        for key, cell in table.get_celld().items():
            cell.set_linewidth(2)
            if key[0] == 0:  # 寻找表头
                cell.set_text_props(fontsize=26, fontweight='bold')  # 设置表头字号为18
                cell.set_facecolor('lightgrey')
            if key[0] != 0:  # 跳过表头
                cell.set_fontsize(24)

        significant = item_StaticResult.select("p_value").to_numpy().flatten() < 0.05
        is_significant = significant.any()
        if is_significant:
            for ax in [ax_box, ax_quantile]:
                for spine in ax.spines.values():
                    spine.set_edgecolor('red')
            for i, sig in enumerate(significant):
                if sig:
                    for j in range(len(table_data[0])):
                        table[(i+1, j)].set_facecolor('orange')

        plt.subplots_adjust(hspace=0.2, wspace=0.01)
        plt.tight_layout()

        # 将图像保存到内存中
        image_stream = BytesIO()
        plt.savefig(image_stream, format='png')
        plt.close(fig)
        image_stream.seek(0)
        return image_stream, is_significant
    except Exception as e:
        print(f"Error in plot_item for item {item}: {e}")
        return None, False  # 确保返回值不是 None

def create_ppt(image_streams, ppt_template_path, output_dir, ppt_name, status_label, root, item_sort_df):
    prs = Presentation(ppt_template_path)
    slide_width, slide_height = prs.slide_width, prs.slide_height
    image_width, image_height = Inches(2.930), Inches(1.8307)
    margin_x, margin_y = (slide_width - image_width * 4) / 2, (slide_height - image_height * 2) / 2

    def add_images_to_slides(image_list, title_text):
        # 按照 item_sort 排序
        image_list.sort(key=lambda x: item_sort_df.filter(pl.col("item") == x[2])["item_sort"][0])

        for i in range(0, len(image_list), 8):
            slide = prs.slides.add_slide(prs.slide_layouts[2])  # 使用带标题的布局
            slide.shapes.title.text = title_text
            for j, (image_stream, _, item) in enumerate(image_list[i:i+8]):
                left, top = margin_x + (j % 4) * image_width, margin_y + (j // 4) * image_height
                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)

    # 分离显著性图片和非显著性图片
    significant_images = [(img_stream, is_significant, item) for img_stream, is_significant, item in image_streams if is_significant]
    non_significant_images = [(img_stream, is_significant, item) for img_stream, is_significant, item in image_streams if not is_significant]

    # 添加显著性图片
    add_images_to_slides(significant_images, "Metrics Mismatch")
    # 添加非显著性图片
    add_images_to_slides(non_significant_images, "Metrics Comparable")

    ppt_output_path = os.path.join(output_dir, ppt_name)
    prs.save(ppt_output_path)
    print(f"PPT saved at {ppt_output_path}")

    # 更新状态为报告生成完成
    status_label.config(text="Congratulation! Report Generated!")
    root.update_idletasks()

    # 启用打开PPT按钮
    open_ppt_button.config(state='normal')

def start_processing(data, csv_path_var, group_col_var, item_col_var, control_group_var, data_structure_var, data_col_listbox, ppt_template_path_var, sample_size_var, progress_var, status_label, root, plot_mode_var):
    selected_data_cols = [data_col_listbox.get(i) for i in data_col_listbox.curselection()]
    if not selected_data_cols:
        messagebox.showwarning("警告", "未选择数据列，程序终止。")
        return

    try:
        sample_size = int(sample_size_var.get())
        if sample_size < 10000:
            messagebox.showerror("错误", "采样大小不能小于10000。")
            return
    except ValueError:
        messagebox.showerror("错误", "采样大小必须是一个整数。")
        return

    # 生成 item_sort 映射表
    item_sort_df = generate_item_sort(data, item_col_var.get(), selected_data_cols, data_structure_var.get())

    StaticResult, item_data_dict = process_data(data, group_col_var.get(), item_col_var.get(), selected_data_cols, data_structure_var.get(), control_group_var.get(), progress_var, status_label, root, csv_path_var.get())
    csv_file_name = os.path.basename(csv_path_var.get())
    csv_dir = os.path.dirname(csv_path_var.get())
    ppt_name = f"{os.path.splitext(csv_file_name)[0]}_Statistic_Report.pptx"
    items = list(item_data_dict.keys())

    # 更新状态为开始生成图片
    status_label.config(text="Plotting...")
    root.update_idletasks()

    image_streams = []
    if plot_mode_var.get() == "single":
        # 单核绘图
        for idx, item in enumerate(items):
            image_stream, is_significant = plot_item(item, StaticResult, control_group_var.get(), item_data_dict, sample_size)
            image_streams.append((image_stream, is_significant, item))  # 添加 item
            # 更新进度条
            progress_var.set(50 + (idx + 1) / len(items) * 50)  # 图片生成占总进度的50%
            root.update_idletasks()  # 强制刷新界面
    else:
        # 多核绘图
        with ProcessPoolExecutor() as executor:
            future_to_item = {executor.submit(plot_item, item, StaticResult, control_group_var.get(), item_data_dict, sample_size): item for item in items}
            for idx, future in enumerate(as_completed(future_to_item)):
                item = future_to_item[future]
                try:
                    image_stream, is_significant = future.result()
                    image_streams.append((image_stream, is_significant, item))  # 添加 item
                    # 更新进度条
                    progress_var.set(50 + (idx + 1) / len(items) * 50)  # 图片生成占总进度的50%
                    root.update_idletasks()  # 强制刷新界面
                except Exception as e:
                    print(f"Error processing item {item}: {e}")

    # 更新状态为图片生成完成
    status_label.config(text="Plot Created")
    root.update_idletasks()

    create_ppt(image_streams, ppt_template_path_var.get(), csv_dir, ppt_name, status_label, root, item_sort_df)
    progress_var.set(100)  # 完成

import polars as pl
from tkinter import filedialog, messagebox

def select_csv_file(csv_path_var, group_col_combobox, item_col_combobox, data_col_listbox, status_label, root):
    file_path = filedialog.askopenfilename(title="Select CSV file", filetypes=[("CSV文件", "*.csv")])
    csv_path_var.set(file_path)
    
    if file_path:
        try:
            # 更新状态为CSV文件加载中
            status_label.config(text="CSV File Loading...")
            root.update_idletasks()

            print("正在加载CSV文件...")
            data = pl.read_csv(file_path, infer_schema_length=10000)

           
            if data.is_empty():
                messagebox.showwarning("警告", "CSV文件为空或仅包含空行，程序终止。")
                return None

            print("CSV文件加载完成。")
            columns = data.columns

            # 更新界面控件
            group_col_combobox['values'] = list(columns)
            item_col_combobox['values'] = list(columns)
            data_col_listbox.delete(0, 'end')
            for col in columns:
                data_col_listbox.insert('end', col)

            # 更新状态为CSV文件加载完成
            status_label.config(text="CSV File Loaded")
            root.update_idletasks()

            return data

        except Exception as e:
            messagebox.showerror("错误", f"读取CSV文件时出错: {e}")
            return None


def update_control_group_options(data, group_col_var, control_group_combobox, control_group_var):
    if group_col_var.get() and data is not None:
        control_group_values = list(data.select(group_col_var.get()).unique().to_series())
        control_group_combobox['values'] = control_group_values
        control_group_var.set(control_group_values[0] if control_group_values else "")

def toggle_item_selection(data_structure_var, item_frame):
    if data_structure_var.get() == "multiple":
        item_frame.grid_remove()
    else:
        item_frame.grid()

def select_ppt_template(ppt_template_path_var, config_manager):
    file_path = filedialog.askopenfilename(title="Select PPT template", filetypes=[("PPT文件", "*.pptx")])
    if file_path:
        ppt_template_path_var.set(file_path)
        # 保存PPT模板路径到配置文件
        config_manager.save_config(file_path)

def run_in_thread(func, *args):
    thread = threading.Thread(target=func, args=args)
    thread.start()

def open_ppt(ppt_path):
    webbrowser.open(ppt_path)

def toggle_item_selection(data_structure_var, metrics_label, item_col_combobox):
    if data_structure_var.get() == "multiple":
        metrics_label.grid_remove()
        item_col_combobox.grid_remove()
    else:
        metrics_label.grid()
        item_col_combobox.grid()

def main():
    global open_ppt_button

    # 创建配置管理器
    config_manager = ConfigManager()

    root = Tk()
    root.title("Chartify Pro")
    root.geometry("425x750")
    root.resizable(False, True)
    data = None
    csv_path_var = StringVar()
    ppt_template_path_var = StringVar()
    sample_size_var = StringVar(value="10000")  # 默认值为10000
    progress_var = DoubleVar(value=0)  # 进度条变量
    plot_mode_var = StringVar(root, value="single")  # 默认单核
    
    # 加载保存的PPT模板路径
    saved_template = config_manager.load_config()
    if saved_template:
        ppt_template_path_var.set(saved_template)

    def on_select_csv_file():
        nonlocal data
        data = select_csv_file(csv_path_var, group_col_combobox, item_col_combobox, data_col_listbox, status_label, root)
        if data is not None:
            update_control_group_options(data, group_col_var, control_group_combobox, control_group_var)

    Button(root, text="Select CSV File", command=on_select_csv_file, width=20).grid(row=0, column=0, padx=10, pady=5, sticky='w')
    Label(root, textvariable=csv_path_var, width=30, anchor='w', wraplength=300).grid(row=0, column=1, padx=10, pady=5, sticky='w')
    Button(root, text="Select PPT Template", command=lambda: select_ppt_template(ppt_template_path_var, config_manager), width=20).grid(row=1, column=0, padx=10, pady=5, sticky='w')
    Label(root, textvariable=ppt_template_path_var, width=30, anchor='w', wraplength=300).grid(row=1, column=1, padx=10, pady=5, sticky='w')

    data_structure_var = StringVar(root, value="single")
    group_col_var = StringVar(root)
    item_col_var = StringVar(root)
    control_group_var = StringVar(root)

    Label(root, text="Select Table Type:", width=20, anchor='center').grid(row=3, column=0, padx=10, pady=5, sticky='w')
    structure_frame = Frame(root)
    structure_frame.grid(row=3, column=1, padx=10, pady=5, sticky='w')
    Radiobutton(structure_frame, text="Single", variable=data_structure_var, value="single", width=10).pack(side='left')
    Radiobutton(structure_frame, text="Multiple", variable=data_structure_var, value="multiple", width=10).pack(side='left')

    # 添加绘图模式选择框
    Label(root, text="Plot Mode:", width=20, anchor='center').grid(row=4, column=0, padx=10, pady=5, sticky='w')
    plot_mode_frame = Frame(root)
    plot_mode_frame.grid(row=4, column=1, padx=10, pady=5, sticky='w')
    Radiobutton(plot_mode_frame, text="Single Core", variable=plot_mode_var, value="single", width=10).pack(side='left')
    Radiobutton(plot_mode_frame, text="Multi Cores", variable=plot_mode_var, value="multi", width=10).pack(side='left')

    Label(root, text="Select Group Column:", width=20, anchor='center').grid(row=5, column=0, padx=10, pady=5, sticky='w')
    group_col_combobox = Combobox(root, textvariable=group_col_var, width=28)
    group_col_combobox.grid(row=5, column=1, padx=10, pady=5, sticky='w')

    Label(root, text="Select Control Group:", width=20, anchor='center').grid(row=6, column=0, padx=10, pady=5, sticky='w')
    control_group_combobox = Combobox(root, textvariable=control_group_var, width=28)
    control_group_combobox.grid(row=6, column=1, padx=10, pady=5, sticky='w')

     # 直接在主窗口中创建控件
    metrics_label = Label(root, text="Select Metrics Column:", width=20, anchor='center')
    metrics_label.grid(row=8, column=0, padx=10, pady=5, sticky='w')
    item_col_combobox = Combobox(root, textvariable=item_col_var, width=28)
    item_col_combobox.grid(row=8, column=1, padx=10, pady=5, sticky='w')

    Label(root, text="Select Data column(s):", width=20, anchor='center').grid(row=9, column=0, padx=10, pady=5, sticky='w')
    data_col_listbox = Listbox(root, selectmode=EXTENDED, width=30)
    data_col_listbox.grid(row=9, column=1, padx=10, pady=5, sticky='w')

    Label(root, text="Sample Size:", width=20, anchor='center').grid(row=10, column=0, padx=10, pady=5, sticky='w')
    sample_size_entry = Entry(root, textvariable=sample_size_var, width=30)
    sample_size_entry.grid(row=10, column=1, padx=10, pady=5, sticky='w')
    
    Button(root, text="Auto Gen Start!", command=lambda: run_in_thread(start_processing, data, csv_path_var, group_col_var, item_col_var, control_group_var, data_structure_var, data_col_listbox, ppt_template_path_var, sample_size_var, progress_var, status_label, root, plot_mode_var), width=20).grid(row=12, column=0, columnspan=2, padx=10, pady=5, sticky='ew')

    # 添加进度条
    progress_bar = ttk.Progressbar(root, variable=progress_var, maximum=100, length=300)
    progress_bar.grid(row=11, column=0, columnspan=2, padx=10, pady=10, sticky='ew')

    # 添加状态标签
    status_label = Label(root, text="", width=40, anchor='center')
    status_label.grid(row=10, column=0, columnspan=2, padx=10, pady=5, sticky='ew')

    data_structure_var.trace_add("write", lambda *args: toggle_item_selection(data_structure_var, metrics_label, item_col_combobox))

    # 添加打开PPT按钮
    open_ppt_button = Button(root, text="Open Report", command=lambda: open_ppt(os.path.join(os.path.dirname(csv_path_var.get()), f"{os.path.splitext(os.path.basename(csv_path_var.get()))[0]}_Statistic_Report.pptx")), width=20, state='disabled')
    open_ppt_button.grid(row=13, column=0, columnspan=2, padx=10, pady=5, sticky='ew')

    # 监听group_col_var的变化以更新控制组选项
    group_col_var.trace_add("write", lambda *args: update_control_group_options(data, group_col_var, control_group_combobox, control_group_var))

    root.mainloop()

if __name__ == "__main__":
    import multiprocessing
    multiprocessing.freeze_support()
    main()